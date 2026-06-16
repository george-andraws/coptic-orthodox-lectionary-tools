#!/usr/bin/env python3
"""Build Phase 6 handoff deck for the Coptic lectionary design layer."""
from __future__ import annotations

import csv
import json
import re
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "presentation"
DESIGN = ROOT / "out" / "design"
OUT.mkdir(parents=True, exist_ok=True)

FORBIDDEN_WORDS = ["delve", "multifaceted", "additionally", "landscape", "underscore", "foster", "interplay"]

COLORS = {
    "burgundy": RGBColor(93, 23, 37),
    "gold": RGBColor(198, 151, 73),
    "cream": RGBColor(248, 242, 230),
    "charcoal": RGBColor(35, 38, 42),
    "muted": RGBColor(97, 94, 87),
    "sage": RGBColor(107, 131, 115),
    "white": RGBColor(255, 255, 255),
}


def read_csv(path: Path) -> list[dict]:
    with path.open(newline="", encoding="utf-8") as f:
        return list(csv.DictReader(f))


def add_bg(slide, color="cream"):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS[color]


def add_text(slide, text, x, y, w, h, size=24, bold=False, color="charcoal", align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = "Aptos"
    p.font.color.rgb = COLORS[color]
    if align:
        p.alignment = align
    return box


def add_title(slide, title, subtitle=None, dark=False):
    add_text(slide, title, 0.7, 0.45, 11.9, 0.65, size=34, bold=True, color="white" if dark else "burgundy")
    if subtitle:
        add_text(slide, subtitle, 0.72, 1.1, 11.5, 0.45, size=16, color="cream" if dark else "muted")


def add_card(slide, title, body, x, y, w, h, accent="gold"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]
    shape.line.color.rgb = COLORS[accent]
    shape.line.width = Pt(1.5)
    add_text(slide, title, x + 0.18, y + 0.15, w - 0.36, 0.35, size=16, bold=True, color="burgundy")
    add_text(slide, body, x + 0.18, y + 0.62, w - 0.36, h - 0.75, size=13, color="charcoal")


def add_stat(slide, number, label, x, y, w=2.2, color="burgundy"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(1.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]
    shape.line.color.rgb = COLORS["gold"]
    add_text(slide, str(number), x + 0.08, y + 0.12, w - 0.16, 0.45, size=24, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.12, y + 0.65, w - 0.24, 0.36, size=10, color="muted", align=PP_ALIGN.CENTER)


def add_footer(slide, n):
    add_text(slide, f"Light and Logos lectionary handoff | {n}", 0.72, 7.15, 11.7, 0.22, size=8, color="muted", align=PP_ALIGN.RIGHT)


def bullet_lines(items: list[str]) -> str:
    return "\n".join(f"• {item}" for item in items)


def build_outline(summary: dict, residue_counts: dict, commem_methods: dict, bridge_conf: dict) -> str:
    outline = f"""# Coptic Lectionary Design Layer Deck Outline

Purpose: handoff deck for George before pushing the article and design-layer data into the site repo.

Visual direction: warm Coptic teaching deck with burgundy, cream, gold, and charcoal. Mostly visual, one idea per slide.

## Slide 1 - The Church teaches Scripture in time
- Title slide.
- Message: the lectionary is not a list. It is Scripture received in worship, season, and commemoration.

## Slide 2 - The problem with date-only tools
- Show date-to-readings, passage-to-uses, and source-status cards.
- Speaker note: the reverse lectionary answers where the Church reads a passage.

## Slide 3 - The design answer is identity first
- Show source label to canonical MT to canonical LXX to identity key.
- Counts: {summary['reading_identity_rows']} identities and {summary['reverse_lectionary_presentation_rows']} presentation rows.

## Slide 4 - Psalm numbering must be honest
- Show Psalm 50 LXX and Psalm 51 MT as one identity with labeled witnesses.
- Speaker note: preserve source labels, do not flatten traditions.

## Slide 5 - Pascha needed attestation, not guessing
- Counts: {summary['pascha_attestation_rows']} Pascha groups, {summary['temporal_residue_rows']} temporal residue rows.
- Speaker note: Coptic Reader fixture governs its captured scope.

## Slide 6 - Temporal classification prevents overclaiming
- Residue counts: {json.dumps(residue_counts, sort_keys=True)}.
- Speaker note: unresolved rows are classified review residue, not hidden failures.

## Slide 7 - The Synaxarium bridge is useful and humble
- Counts: {summary['synaxarium_commemoration_rows']} commemorations, {summary['synaxarium_bridge_rows']} bridge rows.
- Methods: {json.dumps(commem_methods, sort_keys=True)}.
- Bridge confidence: {json.dumps(bridge_conf, sort_keys=True)}.
- Speaker note: all bridge rows are medium-confidence collection-type discovery links.

## Slide 8 - What the site consumes
- Show article, presentation dataset, today's readings, passage footprint, bridge, open questions, and integration spec.

## Slide 9 - Open questions are batched
- Psalm equivalence, 69 collections list, Coptic Reader coverage beyond fixture, prose-lead Synaxarium wording, site corpus joins.

## Slide 10 - George's push path
- Copy files, wire search to identity keys, accept MT and LXX input, join site slugs, verify the plain URL after deploy.
"""
    return outline


def build_deck(summary: dict, residue_counts: dict, commem_methods: dict, bridge_conf: dict) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1
    s = prs.slides.add_slide(blank)
    add_bg(s, "burgundy")
    add_text(s, "The Coptic Lectionary Design Layer", 0.8, 1.2, 11.8, 0.85, size=36, bold=True, color="white")
    add_text(s, "Reverse lectionary, Pascha attestation, Psalm numbering, and Synaxarium bridge", 0.85, 2.05, 11.3, 0.55, size=18, color="cream")
    add_text(s, "Handoff deck for George before the final site push", 0.85, 5.85, 10.5, 0.45, size=15, color="gold")
    add_footer(s, 1)

    # 2
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Date-only tools hide passage meaning", "The reverse lectionary answers where the Church reads a passage.")
    add_card(s, "Date to readings", "Useful for today, but narrow. It answers what is appointed now.", 0.75, 2.0, 3.75, 2.4)
    add_card(s, "Passage to uses", "Shows Holy Week, feasts, hours, saints, and repeated prayer use.", 4.8, 2.0, 3.75, 2.4, "sage")
    add_card(s, "Source status", "Current, historical, inferred, and unresolved rows stay visible.", 8.85, 2.0, 3.75, 2.4)
    add_footer(s, 2)

    # 3
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Identity first, labels second", "The same reading can have different labels across sources.")
    add_stat(s, summary["reading_identity_rows"], "canonical identities", 0.8, 2.0)
    add_stat(s, summary["reverse_lectionary_presentation_rows"], "presentation rows", 3.3, 2.0)
    add_card(s, "Data shape", bullet_lines(["source label", "canonical MT reference", "canonical LXX reference", "stable identity key", "source provenance"]), 6.0, 1.9, 5.8, 2.8, "sage")
    add_text(s, "Result: the site can accept either search tradition and resolve to one identity.", 1.0, 5.45, 11.2, 0.55, size=20, bold=True, color="burgundy", align=PP_ALIGN.CENTER)
    add_footer(s, 3)

    # 4
    s = prs.slides.add_slide(blank)
    add_bg(s, "charcoal")
    add_title(s, "Psalm numbering must be honest", "Preserve both traditions instead of forcing one label.", dark=True)
    add_card(s, "LXX liturgical", "Psalm 50\nCoptic Reader fixture label", 1.0, 2.1, 3.4, 2.2)
    add_card(s, "MT / modern English", "Psalm 51\ncommon English search label", 5.0, 2.1, 3.4, 2.2, "sage")
    add_card(s, "Identity key", "one reading identity\nwith source convention noted", 9.0, 2.1, 3.2, 2.2)
    add_text(s, "False conflicts disappear when numbering convention is explicit.", 1.0, 5.65, 11.5, 0.45, size=20, color="gold", align=PP_ALIGN.CENTER)
    add_footer(s, 4)

    # 5
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Pascha needed attestation, not guessing", "Current practice and historical witnesses are stored separately.")
    add_stat(s, summary["pascha_attestation_rows"], "Pascha attestation groups", 0.9, 2.0)
    add_stat(s, summary["temporal_classification_rows"], "temporal rows", 3.4, 2.0)
    add_stat(s, summary["temporal_residue_rows"], "review residue rows", 5.9, 2.0)
    add_card(s, "Rule", "Coptic Reader governs the captured Wednesday Day fixture. Other rows stay cited and classified, not silently promoted.", 8.6, 1.85, 3.8, 2.6, "sage")
    add_footer(s, 5)

    # 6
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Temporal classification prevents overclaiming", "Unsettled rows are review residue, not hidden failures.")
    x = 0.85
    for label, count in sorted(residue_counts.items()):
        add_stat(s, count, label.replace("_", " ")[:34], x, 2.05, w=2.9)
        x += 3.05
        if x > 10.5:
            x = 0.85
    add_text(s, "The manifest records true source disagreement as zero for this run.", 0.9, 5.85, 11.3, 0.35, size=17, bold=True, color="burgundy", align=PP_ALIGN.CENTER)
    add_footer(s, 6)

    # 7
    s = prs.slides.add_slide(blank)
    add_bg(s, "cream")
    add_title(s, "The Synaxarium bridge is useful and humble", "Discovery links, not direct proper-reading proof.")
    add_stat(s, summary["synaxarium_commemoration_rows"], "commemorations", 0.85, 2.0)
    add_stat(s, summary["synaxarium_bridge_rows"], "bridge rows", 3.35, 2.0)
    add_stat(s, bridge_conf.get("medium", 0), "medium confidence", 5.85, 2.0)
    add_card(s, "Important limit", "All bridge rows are collection-type links. Repeated slots are source-row or variant catalog entries, not a resolved service schedule.", 8.45, 1.85, 3.9, 2.85, "sage")
    add_footer(s, 7)

    # 8
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "The site consumes a push package", "Article, data, integration rules, and open questions are separated.")
    add_card(s, "Reader-facing", "coptic-lectionary-and-synaxarium.md\nsite_integration_spec.md", 0.8, 1.9, 3.8, 2.1)
    add_card(s, "Data", "reverse presentation\ntoday snapshot\npassage footprint\nSynaxarium bridge", 4.8, 1.9, 3.8, 2.1, "sage")
    add_card(s, "Controls", "schema\nsource registry\nPsalm crosswalk\nopen questions", 8.8, 1.9, 3.7, 2.1)
    add_text(s, "George does the final site push and live plain-URL verification.", 1.0, 5.55, 11.2, 0.5, size=20, bold=True, color="burgundy", align=PP_ALIGN.CENTER)
    add_footer(s, 8)

    # 9
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Open questions are batched", "The run did not stop midstream, but it did not hide unsettled evidence.")
    add_card(s, "Needs later review", bullet_lines(["Psalm 41:1 equivalence", "full list of 69 collections", "Coptic Reader beyond fixture", "141 prose-lead Synaxarium titles", "site corpus slug joins"]), 0.9, 1.85, 5.7, 3.55, "gold")
    add_card(s, "Where to look", "audit_artifacts/open_questions_for_george.md\nout/design/temporal_residue.csv\nout/design/synaxarium_commemorations.csv", 7.0, 1.85, 5.1, 3.55, "sage")
    add_footer(s, 9)

    # 10
    s = prs.slides.add_slide(blank)
    add_bg(s, "burgundy")
    add_title(s, "George's push path", "Copy, wire, join, verify.", dark=True)
    add_card(s, "1. Copy files", "Use the file list in site_integration_spec.md", 0.85, 2.0, 2.8, 2.1)
    add_card(s, "2. Wire search", "Accept MT and LXX input, then map to identity_key", 3.85, 2.0, 2.8, 2.1, "sage")
    add_card(s, "3. Join site data", "Fill homily, chapter-study, and audio slugs in coptic-corpus", 6.85, 2.0, 2.8, 2.1)
    add_card(s, "4. Verify live", "Check the plain public URL after deploy", 9.85, 2.0, 2.7, 2.1)
    add_text(s, "Do not overclaim: source status travels with every row.", 1.0, 5.65, 11.3, 0.5, size=20, color="gold", align=PP_ALIGN.CENTER)
    add_footer(s, 10)
    return prs


def collect_text(prs: Presentation) -> str:
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                parts.append(shape.text)
    return "\n".join(parts)


def main() -> None:
    summary = json.loads((DESIGN / "BUILD_DESIGN_SUMMARY.json").read_text(encoding="utf-8"))
    temporal_residue = read_csv(DESIGN / "temporal_residue.csv")
    commems = read_csv(DESIGN / "synaxarium_commemorations.csv")
    bridge = read_csv(DESIGN / "synaxarium_reading_bridge.csv")
    residue_counts = dict(Counter(r["residue_type"] for r in temporal_residue))
    commem_methods = dict(Counter(r["extraction_method"] for r in commems))
    bridge_conf = dict(Counter(r["confidence"] for r in bridge))

    outline = build_outline(summary, residue_counts, commem_methods, bridge_conf)
    for word in FORBIDDEN_WORDS:
        if re.search(rf"\b{re.escape(word)}\b", outline, re.I):
            raise AssertionError(f"Forbidden word in deck outline: {word}")
    if "—" in outline:
        raise AssertionError("Em dash found in deck outline")
    (OUT / "lectionary_design_layer_deck_outline.md").write_text(outline, encoding="utf-8")

    prs = build_deck(summary, residue_counts, commem_methods, bridge_conf)
    text = collect_text(prs)
    for word in FORBIDDEN_WORDS:
        if re.search(rf"\b{re.escape(word)}\b", text, re.I):
            raise AssertionError(f"Forbidden word in deck: {word}")
    if "—" in text:
        raise AssertionError("Em dash found in deck")
    if len(prs.slides) != 10:
        raise AssertionError(f"Expected 10 slides, got {len(prs.slides)}")
    prs.save(OUT / "lectionary_design_layer_deck.pptx")
    print(json.dumps({
        "pptx": str(OUT / "lectionary_design_layer_deck.pptx"),
        "outline": str(OUT / "lectionary_design_layer_deck_outline.md"),
        "slides": len(prs.slides),
    }, indent=2))


if __name__ == "__main__":
    main()
